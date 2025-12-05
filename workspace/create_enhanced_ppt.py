#!/usr/bin/env python3
"""
使用 python-pptx 创建包含图片和完整样式的 PPT 演示文稿
"""

import os
import sys
from pathlib import Path
from PIL import Image

def create_enhanced_ppt():
    """创建包含图片和完整样式的增强版 PPT"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.dml import MSO_THEME_COLOR
        print("正在创建增强版 PowerPoint 演示文稿...")
    except ImportError:
        print("❌ python-pptx 库未安装")
        print("💡 您可以运行: pip install python-pptx")
        return False

    try:
        # 创建演示文稿
        prs = Presentation()

        # 设置演示文稿尺寸为16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        current_dir = Path(__file__).parent
        images_dir = current_dir / "images"

        # 幻灯片数据
        slides_data = [
            {
                "title": "AI辅助编程进入巡航模式",
                "content": "三种主流AI编程工具的实践对比与经验分享",
                "layout": "title_only",
                "background_color": (28, 40, 51)  # #1C2833
            },
            {
                "title": "Vue-CopilotKit 升级项目",
                "content": ["主线需求：Vue-CopilotKit 升级到与 React 新版兼容",
                           "支线需求：后端接口优化，支持多模型切换"],
                "layout": "two_content",
                "image": "20251205105230499.png",
                "background_color": (28, 40, 51)
            },
            {
                "title": "三种AI编程方案对比",
                "content": [
                    "1. Claude Code + GLM4.6：1-2周，60元/月",
                    "2. Antigravity + Gemini3 Pro：4小时，按需计费",
                    "3. Auto-coder + Deepseek3.2：2天，24.59元"
                ],
                "layout": "title_and_content",
                "background_color": (28, 40, 51)
            },
            {
                "title": "方案实施结果",
                "content": ["项目成功升级到1.10.6+", "与React新版完全兼容", "后端接口优化完成"],
                "layout": "comparison",
                "images": ["20251205114008211.png", "20251205114207002.png", "20251205114434093.png"],
                "background_color": (28, 40, 51)
            },
            {
                "title": "实施方法：项目设计",
                "content": ["目录结构三目录设计", "配置文档指导", "AI工具正确配置"],
                "layout": "two_content",
                "image": "20251205114900758.png",
                "background_color": (28, 40, 51)
            },
            {
                "title": "配置对比与总结",
                "content": ["CLAUDE.md vs autocoder RULES.md", "Antigrativity 配置优势", "核心洞察：好的开始是成功的一半"],
                "layout": "title_and_content",
                "background_color": (28, 40, 51)
            }
        ]

        def set_background(slide, color):
            """设置幻灯片背景颜色"""
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*color)

        def create_title_slide(title, subtitle, bg_color):
            """创建标题幻灯片"""
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片
            slide = prs.slides.add_slide(slide_layout)

            set_background(slide, bg_color)

            # 设置标题
            title_shape = slide.shapes.title
            title_shape.text = title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER

            # 设置副标题
            if slide.placeholders:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle
                for paragraph in subtitle_shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(28)
                    paragraph.font.color.rgb = RGBColor(243, 156, 18)  # #F39C12
                    paragraph.alignment = PP_ALIGN.CENTER

        def create_content_slide(title, content, layout, image_path=None, images=None, bg_color=(28, 40, 51)):
            """创建内容幻灯片"""
            if layout == "title_only":
                slide_layout = prs.slide_layouts[5]  # 标题
                slide = prs.slides.add_slide(slide_layout)
            else:
                slide_layout = prs.slide_layouts[1]  # 标题和内容
                slide = prs.slides.add_slide(slide_layout)

            set_background(slide, bg_color)

            # 设置标题
            title_shape = slide.shapes.title
            title_shape.text = title
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # 添加内容
            if isinstance(content, list):
                content_text = '\n'.join(content)
            else:
                content_text = content

            if slide.placeholders:
                content_shape = slide.placeholders[1] if layout != "title_only" else None
                if content_shape:
                    content_shape.text = content_text
                    for paragraph in content_shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(244, 246, 246)

            # 添加图片
            if image_path and images_dir:
                full_image_path = images_dir / image_path
                if full_image_path.exists():
                    try:
                        # 调整图片大小
                        img = Image.open(full_image_path)
                        img_width, img_height = img.size

                        # 计算适当的显示大小
                        max_width = Inches(5)
                        max_height = Inches(3.5)

                        if img_width > img_height:
                            width = max_width
                            height = img_height * (max_width / img_width)
                        else:
                            height = max_height
                            width = img_width * (max_height / img_height)

                        # 添加图片到幻灯片
                        left = Inches(5.5)
                        top = Inches(2)
                        slide.shapes.add_picture(str(full_image_path), left, top, width=width, height=height)
                    except Exception as e:
                        print(f"⚠️  无法添加图片 {image_path}: {e}")

            # 添加多张图片（对比布局）
            if images and images_dir:
                for i, img_name in enumerate(images):
                    full_image_path = images_dir / img_name
                    if full_image_path.exists():
                        try:
                            left = Inches(0.5 + i * 3.2)
                            top = Inches(2.5)
                            width = Inches(3)
                            slide.shapes.add_picture(str(full_image_path), left, top, width=width)
                        except Exception as e:
                            print(f"⚠️  无法添加图片 {img_name}: {e}")

        # 创建幻灯片
        for i, slide_data in enumerate(slides_data):
            print(f"正在创建第 {i+1} 张幻灯片: {slide_data['title']}")

            if i == 0:  # 第一张幻灯片使用标题布局
                create_title_slide(
                    slide_data["title"],
                    slide_data["content"],
                    slide_data["background_color"]
                )
            else:
                create_content_slide(
                    slide_data["title"],
                    slide_data["content"],
                    slide_data["layout"],
                    slide_data.get("image"),
                    slide_data.get("images"),
                    slide_data["background_color"]
                )

        # 保存文件
        output_path = current_dir / 'AI辅助编程进入巡航模式_增强版.pptx'
        prs.save(output_path)

        print(f"✅ 增强版 PowerPoint 演示文稿创建成功!")
        print(f"📁 保存位置: {output_path}")
        print(f"📊 包含 {len(slides_data)} 张幻灯片")
        return True

    except Exception as e:
        print(f"❌ 创建增强版 PPT 时发生错误: {e}")
        return False

if __name__ == "__main__":
    print("🚀 开始创建增强版演示文稿...")

    if create_enhanced_ppt():
        print("\n✨ 任务完成!")
    else:
        print("\n❌ 任务失败，请检查依赖库安装")