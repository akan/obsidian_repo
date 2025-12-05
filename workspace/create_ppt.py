#!/usr/bin/env python3
"""
使用 Python 创建 PPT 演示文稿
如果 python-pptx 不可用，则创建 HTML 版本作为替代
"""

import os
import sys
from pathlib import Path

def create_html_presentation():
    """创建 HTML 版本的演示文稿作为备选方案"""

    print("正在创建 HTML 演示文稿...")

    # 读取现有的 HTML 幻灯片文件
    slide_files = [
        'slide1.html',
        'slide2.html',
        'slide3.html',
        'slide4.html',
        'slide5.html',
        'slide6.html'
    ]

    # 创建合并的 HTML 演示文稿
    html_content = '''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>AI辅助编程进入巡航模式 - 演示文稿</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        body {
            font-family: 'Inter', Arial, sans-serif;
            background: #1C2833;
            color: white;
            overflow-x: hidden;
        }

        .slide {
            width: 720pt;
            height: 405pt;
            margin: 20pt auto;
            background: linear-gradient(135deg, #1C2833 0%, #2E4053 100%);
            border-radius: 12pt;
            overflow: hidden;
            position: relative;
            box-shadow: 0 10pt 30pt rgba(0,0,0,0.3);
        }

        .slide-number {
            position: absolute;
            bottom: 20pt;
            right: 30pt;
            font-size: 14pt;
            color: #AAB7B8;
            background: rgba(0,0,0,0.3);
            padding: 5pt 12pt;
            border-radius: 20pt;
        }

        .navigation {
            text-align: center;
            margin: 30pt 0;
        }

        .nav-btn {
            background: #F39C12;
            color: #1C2833;
            border: none;
            padding: 12pt 24pt;
            margin: 0 10pt;
            border-radius: 25pt;
            font-size: 14pt;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
        }

        .nav-btn:hover {
            background: #E67E22;
            transform: translateY(-2pt);
        }

        .nav-btn:disabled {
            background: #7F8C8D;
            cursor: not-allowed;
            transform: none;
        }

        .slide-info {
            text-align: center;
            margin: 20pt 0;
            color: #AAB7B8;
            font-size: 16pt;
        }

        @media print {
            body { background: white; }
            .slide {
                margin: 0;
                page-break-after: always;
                box-shadow: none;
            }
            .navigation, .slide-info { display: none; }
        }
    </style>
</head>
<body>
    <div class="slide-info">
        <strong>AI辅助编程进入巡航模式</strong> - 演示文稿
    </div>

    <div id="slide-container">
'''

    current_dir = Path(__file__).parent

    # 读取并整合所有幻灯片
    for i, slide_file in enumerate(slide_files):
        slide_path = current_dir / slide_file
        if slide_path.exists():
            with open(slide_path, 'r', encoding='utf-8') as f:
                slide_html = f.read()
                # 提取 body 内容
                if '<body>' in slide_html and '</body>' in slide_html:
                    body_start = slide_html.find('<body>') + 6
                    body_end = slide_html.find('</body>')
                    body_content = slide_html[body_start:body_end]

                    html_content += f'''
        <div class="slide" id="slide-{i+1}">
            {body_content}
            <div class="slide-number">{i+1} / {len(slide_files)}</div>
        </div>
'''

    html_content += '''
    </div>

    <div class="navigation">
        <button class="nav-btn" onclick="previousSlide()">上一张</button>
        <button class="nav-btn" onclick="nextSlide()">下一张</button>
        <button class="nav-btn" onclick="toggleFullscreen()">全屏</button>
    </div>

    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');

        function showSlide(index) {
            slides.forEach(slide => slide.style.display = 'none');
            if (index >= 0 && index < slides.length) {
                slides[index].style.display = 'block';
                currentSlide = index;
                updateNavigation();
            }
        }

        function nextSlide() {
            if (currentSlide < slides.length - 1) {
                showSlide(currentSlide + 1);
            }
        }

        function previousSlide() {
            if (currentSlide > 0) {
                showSlide(currentSlide - 1);
            }
        }

        function updateNavigation() {
            const prevBtn = document.querySelector('.nav-btn');
            const nextBtn = document.querySelectorAll('.nav-btn')[1];

            prevBtn.disabled = currentSlide === 0;
            nextBtn.disabled = currentSlide === slides.length - 1;
        }

        function toggleFullscreen() {
            if (!document.fullscreenElement) {
                document.documentElement.requestFullscreen();
            } else {
                document.exitFullscreen();
            }
        }

        // 键盘导航
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight' || e.key === ' ') nextSlide();
            if (e.key === 'ArrowLeft') previousSlide();
            if (e.key === 'f' || e.key === 'F') toggleFullscreen();
        });

        // 初始化
        showSlide(0);
    </script>
</body>
</html>'''

    # 保存 HTML 文件
    output_path = current_dir / 'AI辅助编程进入巡航模式.html'
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)

    print(f"✅ HTML 演示文稿创建成功!")
    print(f"📁 保存位置: {output_path}")
    print(f"🌐 在浏览器中打开此文件即可查看演示文稿")
    print(f"⌨️  使用键盘方向键或按钮导航，按 F 键全屏")

def try_create_pptx():
    """尝试使用 python-pptx 创建 PowerPoint 文件"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        print("正在创建 PowerPoint 演示文稿...")

        # 创建演示文稿
        prs = Presentation()

        # 幻灯片内容
        slides_data = [
            {
                "title": "AI辅助编程进入巡航模式",
                "content": "三种主流AI编程工具的实践对比与经验分享"
            },
            {
                "title": "项目背景",
                "content": "Vue-CopilotKit 升级项目\n从 1.0.1 升级到 1.10.6+\n后端接口优化与功能扩展"
            },
            {
                "title": "三种AI编程方案对比",
                "content": "1. Claude Code + GLM4.6\n2. Antigravity + Gemini3 Pro + Sonnet4.5\n3. Auto-coder + Deepseek3.2"
            },
            {
                "title": "实施方法：项目设计",
                "content": "目录结构设计理念\n• 充足上下文\n• 明确参考方案\n• Debug支持"
            },
            {
                "title": "配置文档对比",
                "content": "CLAUDE.md vs autocoder RULES.md vs Antigravity 配置\n不同AI工具的配置策略对比"
            },
            {
                "title": "总结与展望",
                "content": "好的开始是成功的一半\n规划比执行更重要\n工具选择决定效率"
            }
        ]

        # 创建幻灯片
        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[1]  # 标题和内容布局
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            title = slide.shapes.title
            title.text = slide_data["title"]

            # 设置内容
            content = slide.placeholders[1]
            content.text = slide_data["content"]

            # 设置标题字体
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER

        # 保存文件
        output_path = Path(__file__).parent / 'AI辅助编程进入巡航模式.pptx'
        prs.save(output_path)

        print(f"✅ PowerPoint 演示文稿创建成功!")
        print(f"📁 保存位置: {output_path}")
        return True

    except ImportError:
        print("⚠️  python-pptx 库未安装")
        print("💡 您可以运行: pip install python-pptx")
        return False
    except Exception as e:
        print(f"❌ 创建 PowerPoint 时发生错误: {e}")
        return False

if __name__ == "__main__":
    print("🚀 开始创建演示文稿...")

    # 首先尝试创建 PPTX
    if not try_create_pptx():
        print("\n🔄 转为创建 HTML 版本...")
        create_html_presentation()

    print("\n✨ 任务完成!")